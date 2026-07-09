import os
import pandas as pd
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders import ( PyPDFLoader, Docx2txtLoader,TextLoader,CSVLoader)


class LightweightExcelLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        df = pd.read_excel(self.file_path)
        content = df.to_string(index=False)
        return [Document(page_content=content, metadata={"source": self.file_path})]

class LightweightPPTXLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        prs = Presentation(self.file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        
        full_text = "\n".join(text_content)
        return [Document(page_content=full_text, metadata={"source": self.file_path})]


def load_document(file_path: str):
    extension = os.path.splitext(file_path)[1].lower()
    try:
        if extension == ".pdf":
            loader = PyPDFLoader(file_path)

        elif extension == ".docx":
            loader = Docx2txtLoader(file_path)

        elif extension == ".txt":
            loader = TextLoader(file_path, encoding="utf-8")

        elif extension == ".csv":
            loader = CSVLoader(file_path)

        elif extension == ".pptx":
            loader = LightweightPPTXLoader(file_path)

        elif extension == ".xlsx":
            loader = LightweightExcelLoader(file_path)

        else:
            raise Exception(f"Unsupported file type: {extension}")

        documents = loader.load()

        # update Metadata
        filename = os.path.basename(file_path)
        for document in documents:
            document.metadata.setdefault("file_name", filename)
            document.metadata.setdefault("extension", extension)
            
        return documents

    except Exception as error:
        raise Exception(f"Failed to load document '{file_path}': {error}")


def load_multiple_documents(file_paths: list[str]):
    documents = []
    try:
        for file_path in file_paths:
            documents.extend(load_document(file_path))

        return documents

    except Exception as error:
        raise Exception(f"Failed to load uploaded documents: {error}")